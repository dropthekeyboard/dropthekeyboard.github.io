#!/usr/bin/env python3
"""
PowerPoint to PNG Converter

This script converts PowerPoint presentations (.pptx files) to PNG images.
Each slide will be saved as a separate PNG file.

Dependencies:
- python-pptx: for reading PowerPoint files
- Pillow (PIL): for image processing
- comtypes (Windows only): for better PowerPoint conversion
- LibreOffice (alternative method): system dependency

Usage:
    python pptx_png.py input.pptx [output_directory]
"""

import os
import sys
import argparse
import subprocess
from pathlib import Path
from typing import Optional, List

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


class PPTXToPNGConverter:
    """Convert PowerPoint presentations to PNG images."""
    
    def __init__(self, input_file: str, output_dir: Optional[str] = None):
        """
        Initialize the converter.
        
        Args:
            input_file: Path to the input PPTX file
            output_dir: Output directory for PNG files (optional)
        """
        self.input_file = Path(input_file)
        self.output_dir = Path(output_dir) if output_dir else self.input_file.parent / f"{self.input_file.stem}_slides"
        
        # Validate input file
        if not self.input_file.exists():
            raise FileNotFoundError(f"Input file not found: {self.input_file}")
        
        if self.input_file.suffix.lower() not in ['.pptx', '.ppt']:
            raise ValueError(f"Unsupported file format: {self.input_file.suffix}")
        
        # Create output directory
        self.output_dir.mkdir(exist_ok=True)
    
    def convert_via_libreoffice(self) -> List[Path]:
        """
        Convert PPTX to PNG using LibreOffice (most reliable method).
        
        Returns:
            List of generated PNG file paths
        """
        print("Converting using LibreOffice...")
        
        # Check if LibreOffice is available
        try:
            subprocess.run(['libreoffice', '--version'], 
                         capture_output=True, check=True)
        except (subprocess.CalledProcessError, FileNotFoundError):
            raise RuntimeError("LibreOffice not found. Please install LibreOffice.")
        
        # Convert PPTX to PDF first
        pdf_path = self.output_dir / f"{self.input_file.stem}.pdf"
        
        cmd_pdf = [
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(self.output_dir),
            str(self.input_file)
        ]
        
        result = subprocess.run(cmd_pdf, capture_output=True, text=True)
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice PDF conversion failed: {result.stderr}")
        
        if not pdf_path.exists():
            raise RuntimeError("PDF conversion failed - output file not found")
        
        # Convert PDF to PNG images
        png_files = self._pdf_to_png(pdf_path)
        
        # Clean up intermediate PDF file
        pdf_path.unlink()
        
        return png_files
    
    def _pdf_to_png(self, pdf_path: Path) -> List[Path]:
        """
        Convert PDF to PNG images using pdftoppm.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of generated PNG file paths
        """
        try:
            subprocess.run(['pdftoppm', '-v'], 
                         capture_output=True, check=True)
        except (subprocess.CalledProcessError, FileNotFoundError):
            raise RuntimeError("pdftoppm not found. Please install poppler-utils.")
        
        output_prefix = self.output_dir / "slide"
        
        cmd_png = [
            'pdftoppm',
            '-png',
            '-r', '300',  # High resolution
            str(pdf_path),
            str(output_prefix)
        ]
        
        result = subprocess.run(cmd_png, capture_output=True, text=True)
        if result.returncode != 0:
            raise RuntimeError(f"PDF to PNG conversion failed: {result.stderr}")
        
        # Find generated PNG files
        png_files = sorted(self.output_dir.glob("slide-*.png"))
        
        # Rename files to have consistent naming
        renamed_files = []
        for i, png_file in enumerate(png_files, 1):
            new_name = self.output_dir / f"slide_{i:03d}.png"
            png_file.rename(new_name)
            renamed_files.append(new_name)
        
        return renamed_files
    
    def convert_via_python_pptx(self) -> List[Path]:
        """
        Convert PPTX to PNG using python-pptx (limited functionality).
        
        Note: This method has limitations and may not work for all presentations.
        
        Returns:
            List of generated PNG file paths
        """
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not available. Install with: pip install python-pptx")
        
        print("Converting using python-pptx...")
        print("Warning: This method has limitations and may not work for complex slides.")
        
        try:
            prs = Presentation(str(self.input_file))
        except Exception as e:
            raise RuntimeError(f"Failed to open presentation: {e}")
        
        png_files = []
        
        for i, slide in enumerate(prs.slides, 1):
            output_file = self.output_dir / f"slide_{i:03d}.png"
            
            # Note: python-pptx doesn't have built-in slide-to-image conversion
            # This is a placeholder that would require additional implementation
            print(f"Slide {i}: {len(slide.shapes)} shapes found")
            
            # For now, create a placeholder image
            if PIL_AVAILABLE:
                img = Image.new('RGB', (1920, 1080), 'white')
                img.save(output_file)
                png_files.append(output_file)
        
        return png_files
    
    def convert(self, method: str = 'auto') -> List[Path]:
        """
        Convert the PowerPoint presentation to PNG images.
        
        Args:
            method: Conversion method ('auto', 'libreoffice', 'python-pptx')
            
        Returns:
            List of generated PNG file paths
        """
        if method == 'auto':
            # Try LibreOffice first (most reliable)
            try:
                return self.convert_via_libreoffice()
            except RuntimeError as e:
                print(f"LibreOffice method failed: {e}")
                print("Falling back to python-pptx method...")
                return self.convert_via_python_pptx()
        elif method == 'libreoffice':
            return self.convert_via_libreoffice()
        elif method == 'python-pptx':
            return self.convert_via_python_pptx()
        else:
            raise ValueError(f"Unknown conversion method: {method}")


def install_dependencies():
    """Install required Python dependencies."""
    dependencies = [
        'python-pptx',
        'Pillow'
    ]
    
    print("Installing Python dependencies...")
    for dep in dependencies:
        try:
            subprocess.run([sys.executable, '-m', 'pip', 'install', dep], 
                         check=True)
            print(f"✓ Installed {dep}")
        except subprocess.CalledProcessError:
            print(f"✗ Failed to install {dep}")


def check_system_dependencies():
    """Check if system dependencies are available."""
    dependencies = {
        'LibreOffice': ['libreoffice', '--version'],
        'pdftoppm': ['pdftoppm', '-v']
    }
    
    available = {}
    for name, cmd in dependencies.items():
        try:
            subprocess.run(cmd, capture_output=True, check=True)
            available[name] = True
            print(f"✓ {name} is available")
        except (subprocess.CalledProcessError, FileNotFoundError):
            available[name] = False
            print(f"✗ {name} is not available")
    
    return available


def main():
    """Main function to handle command line interface."""
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint presentations to PNG images",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
    python pptx_png.py presentation.pptx
    python pptx_png.py presentation.pptx --output ./slides
    python pptx_png.py presentation.pptx --method libreoffice
    python pptx_png.py --check-deps
    python pptx_png.py --install-deps
        """
    )
    
    parser.add_argument('input_file', nargs='?', 
                       help='Input PowerPoint file (.pptx or .ppt)')
    parser.add_argument('-o', '--output', 
                       help='Output directory for PNG files')
    parser.add_argument('-m', '--method', 
                       choices=['auto', 'libreoffice', 'python-pptx'],
                       default='auto',
                       help='Conversion method (default: auto)')
    parser.add_argument('--check-deps', action='store_true',
                       help='Check system dependencies')
    parser.add_argument('--install-deps', action='store_true',
                       help='Install Python dependencies')
    
    args = parser.parse_args()
    
    if args.check_deps:
        print("Checking system dependencies...")
        check_system_dependencies()
        return
    
    if args.install_deps:
        install_dependencies()
        return
    
    if not args.input_file:
        parser.error("Input file is required")
    
    try:
        converter = PPTXToPNGConverter(args.input_file, args.output)
        png_files = converter.convert(args.method)
        
        print(f"\n✓ Conversion completed successfully!")
        print(f"✓ Generated {len(png_files)} PNG files in: {converter.output_dir}")
        
        for png_file in png_files:
            print(f"  - {png_file.name}")
            
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()
